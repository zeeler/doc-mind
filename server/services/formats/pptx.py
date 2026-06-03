"""PPTX 解析 — 逐页提取文本。"""

from pathlib import Path


def parse_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        slide_parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_parts.append(text)
        if slide_parts:
            parts.append(f"## Page {i}")
            parts.append("\n".join(slide_parts))
            parts.append("")
    return "\n".join(parts)
